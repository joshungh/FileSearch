import os
import xlrd
from docx import Document
from pptx import Presentation
from multiprocessing import Process

results = []

def get_drives():
    response = os.popen("wmic logicaldisk get caption")
    list1 = []
    for line in response.readlines():
        line = line.strip("\n")
        line = line.strip("\r")
        line = line.strip(" ")
        if (line == "Caption" or line == ""):
            continue
        list1.append(line + '/')
    return list1

all_drives = get_drives()

def createThreads(value):
    thread_list = []
    for each in all_drives:
        print('drive ' + each + ' Thread Generation completed')
        thread = Process(target=designateThreads, args=(each, value))
        thread_list.append(thread)
    for thread in thread_list:
        thread.start()

def designateThreads(drive, value):
    print('new thread started')
    thread_list = []
    for each in os.listdir(drive):
        print(drive + each)
        if os.path.isdir(drive + each):
            print('starting thread for ' + drive + each)
            thread = Process(target=search_directory, args=(each, drive, value))
            thread_list.append(thread)
        else:
            search_specific(each, drive, value)
    for thread in thread_list:
        thread.start()

def search_directory(each, drive, keyword):
    print('thread search in ' + drive + each + ' started')
    for root, dir, files in os.walk(drive + each, topdown=True):
        for f in files:
            print('search in thread commencing for ' + os.path.join(root, f))
            if os.path.splitext(f)[1] == '.txt':
                with open(root + '/' + f, errors='ignore') as cur_f:
                    if keyword in cur_f.read():
                        results.append(os.path.join(root, f))
                        break

            if os.path.splitext(f)[1] == '.xlsx':
                if f[0] != '~' and f[0] != '$':
                    # wb = xlrd.open_workbook(os.path.expanduser('C:/a1/a1.xlsx'))
                    wb = xlrd.open_workbook(os.path.join(root, f))
                    sheet = wb.sheet_by_index(0)
                    for row_num in range(sheet.nrows):
                        for col_num in range(sheet.ncols):
                            cell_obj = sheet.row(row_num)[col_num]
                            if keyword == cell_obj.value:
                                print(f)
                                results.append(os.path.join(root, f))
                                break

            if os.path.splitext(f)[1] == '.pptx':
                # f = open(os.path.expanduser('~/.' + f)
                if f[0] != '0':
                    prs = Presentation(os.path.join(root, f))
                    for slides in prs.slides:
                        for shape in slides.shapes:
                            if shape.has_text_frame:
                                if (shape.text.find(keyword)) != -1:
                                    results.append(os.path.join(root, f))
                                    break

            if os.path.splitext(f)[1] == '.docx':
                if f[0] != '$' and f[0] != '0' and f[0] != '~':
                    document = Document(root + '/' + f)
                    for p in document.paragraphs:
                        if p.text.find(keyword) != -1:
                            results.append(os.path.join(root, f))
                            break

def search_specific(f, root, keyword):
    print('specific search started on ' + f)
    if os.path.splitext(f)[1] == '.txt':
        if check(f, root, keyword):
            results.append(os.path.join(root, f))

def check(checkFile, filePath, keyword):
    datafile = checkFile
    with open(filePath + '/' + datafile, errors='ignore') as f:
        if keyword in f.read():
            print('Success')
            return True
    return False

if __name__ == '__main__':
    drives = get_drives()
    createThreads('test')